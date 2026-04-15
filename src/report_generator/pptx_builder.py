"""
report_generator/pptx_builder.py — BDA Daily Report PowerPoint Generator.

Architecture: loads assets/template_v1.pptx, modifies each slide in-place.
  • Slide 1  (Cover)    — updates date text only; all branding from template
  • Slide 2  (Agenda)   — 100 % static, untouched
  • Slides 3-11         — cleared and rebuilt with live data

Slide structure (11 slides):
  1.  Cover
  2.  Agenda
  3.  Sumário Executivo
  4.  Liquidez – Moeda Nacional (1/2)  — liquidity table + LUIBOR
  5.  Liquidez – Moeda Nacional (2/2)  — cash-flow + P&L
  6.  Liquidez – Moeda Estrangeira
  7.  Mercado Cambial
  8.  Mercado de Capitais – BODIVA     — segments + stocks
  9.  Mercado de Capitais – Operações BDA
  10. Informação de Mercados (1/2)     — indices + crypto
  11. Informação de Mercados (2/2)     — commodities + minerals

Usage:
    from src.report_generator.pptx_builder import BDAReportGenerator
    gen  = BDAReportGenerator(data)
    path = gen.build("output/bda_report_2026-03-30.pptx")
"""
from __future__ import annotations

import os
import re
import shutil
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Asset paths ────────────────────────────────────────────────────────────────
_ASSETS_DIR  = Path(__file__).parent.parent.parent / "assets"
_TEMPLATE    = _ASSETS_DIR / "template_v1.pptx"
_IMAGES_DIR  = _ASSETS_DIR / "images"

IMG_FOOTER_BANNER = str(_IMAGES_DIR / "footer_banner.png")
IMG_COVER_HEADER  = str(_IMAGES_DIR / "cover_header.png")
IMG_INNER_HEADER  = str(_IMAGES_DIR / "inner_header.png")
IMG_AGENDA_BG     = str(_IMAGES_DIR / "agenda_bg.jpg")

IMG_ICON_S3_LMN       = str(_IMAGES_DIR / "icon_s3_liquidez_mn.png")
IMG_ICON_S3_LME       = str(_IMAGES_DIR / "icon_s3_liquidez_me.png")
IMG_ICON_S3_RENTA     = str(_IMAGES_DIR / "icon_s3_rentabilidade.png")
IMG_ICON_S3_VBAR      = str(_IMAGES_DIR / "icon_s3_vbar.png")
IMG_ICON_S3_CAMBIAL   = str(_IMAGES_DIR / "icon_s3_posicao_cambial.png")
IMG_ICON_S3_REEMB     = str(_IMAGES_DIR / "icon_s3_reembolsos.png")
IMG_ICON_S3_REPORT    = str(_IMAGES_DIR / "icon_report_color.png")
IMG_ICON_CALCULATOR   = str(_IMAGES_DIR / "icon_calculator.png")
IMG_ICON_GEAR         = str(_IMAGES_DIR / "icon_gear.png")
IMG_ICON_KZ_CIRCLE    = str(_IMAGES_DIR / "icon_kz_circle.png")
IMG_ICON_REPORT_MONEY = str(_IMAGES_DIR / "icon_report_money.png")
IMG_ICON_FX_EXCHANGE  = str(_IMAGES_DIR / "icon_fx_exchange.png")
IMG_ICON_GLOBE        = str(_IMAGES_DIR / "icon_globe.png")
IMG_ICON_REPORT_BW    = str(_IMAGES_DIR / "icon_report_bw.png")

# ── Colour palette ─────────────────────────────────────────────────────────────
ORANGE_PRIMARY = RGBColor(0xE8, 0x75, 0x1A)
ORANGE_DARK    = RGBColor(0xC0, 0x50, 0x00)
ORANGE_LIGHT   = RGBColor(0xFF, 0xF0, 0xE0)
BROWN_KPI      = RGBColor(0x5C, 0x2D, 0x00)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
BLACK          = RGBColor(0x00, 0x00, 0x00)
LIGHT_GREY     = RGBColor(0xF5, 0xF5, 0xF5)
MID_GREY       = RGBColor(0xD9, 0xD9, 0xD9)
DARK_GREY      = RGBColor(0x55, 0x55, 0x55)
GREEN_UP       = RGBColor(0x00, 0x7A, 0x33)
RED_DOWN       = RGBColor(0xCC, 0x00, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FONT    = "Calibri"

# ── XML tag set for clearing a slide ──────────────────────────────────────────
_CLEAR_TAGS = {
    qn("p:sp"), qn("p:grpSp"), qn("p:pic"),
    qn("p:graphicFrame"), qn("p:cxnSp"),
}


# ─────────────────────────────────────────────────────────────────────────────
# Slide-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def _clear_slide(slide) -> None:
    """Remove every shape element from a slide's shape-tree."""
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        if child.tag in _CLEAR_TAGS:
            sp_tree.remove(child)


def _update_cover_date(slide, date_str: str) -> None:
    """
    Slide 1 only: find shape_id=13 and update the paragraph that contains
    the date (matches DD.MM.YYYY pattern).
    """
    for shape in slide.shapes:
        if shape.shape_id == 13 and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                if re.match(r"\d{2}\.\d{2}\.\d{4}", para.text.strip()):
                    if para.runs:
                        para.runs[0].text = date_str
                    break
            break


# ─────────────────────────────────────────────────────────────────────────────
# Low-level shape builders
# ─────────────────────────────────────────────────────────────────────────────

def _place_img(slide, path, left, top, width=None, height=None):
    if os.path.isfile(path):
        slide.shapes.add_picture(path, left, top, width, height)


def _add_rect(slide, left, top, width, height,
              fill_color=None, line_color=None, line_width_pt: float = 0.5):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 10, bold: bool = False,
                  italic: bool = False, color: RGBColor = BLACK,
                  align=PP_ALIGN.LEFT, word_wrap: bool = True,
                  font_name: str = FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font_name
    return txBox


# ─────────────────────────────────────────────────────────────────────────────
# Reusable slide furniture
# ─────────────────────────────────────────────────────────────────────────────

def _slide_title(slide, title: str):
    """Inner slide slim header — orange accent block + bold title + orange rule."""
    _add_rect(slide, Inches(0), Inches(0.401), Inches(0.277), Inches(0.397), ORANGE_PRIMARY)
    _add_text_box(slide, title,
                  Inches(0.396), Inches(0.224), Inches(11.321), Inches(0.351),
                  font_size=20, bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    _add_rect(slide, Inches(0.396), Inches(0.607), SLIDE_W - Inches(0.6), Inches(0.015),
              ORANGE_PRIMARY)


def _section_bar(slide, label: str, left, top, width=None, height=Inches(0.28)):
    w = width or (SLIDE_W - left - Inches(0.25))
    _add_rect(slide, left, top, w, height, ORANGE_DARK)
    _add_text_box(slide, label, left + Pt(4), top + Pt(1), w - Pt(8), height - Pt(2),
                  font_size=8, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def _footer(slide):
    footer_h   = Inches(0.72)
    footer_top = SLIDE_H - footer_h
    if os.path.isfile(IMG_FOOTER_BANNER):
        slide.shapes.add_picture(IMG_FOOTER_BANNER, Inches(0), footer_top, SLIDE_W, footer_h)
    else:
        _add_rect(slide, Inches(0), footer_top, SLIDE_W, footer_h, ORANGE_PRIMARY)
        _add_text_box(
            slide,
            "Edifício BDA, Condomínio Dolce Vita, Via S8, Talatona  |  Luanda – Angola",
            Inches(0.3), footer_top + Pt(2), SLIDE_W - Inches(3.5), footer_h - Pt(4),
            font_size=7, color=WHITE, align=PP_ALIGN.LEFT,
        )
        _add_text_box(
            slide, "UMA VISÃO DE FUTURO",
            SLIDE_W - Inches(3.2), footer_top + Pt(2), Inches(3.0), footer_h - Pt(4),
            font_size=7, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
        )


# ─────────────────────────────────────────────────────────────────────────────
# Table helpers
# ─────────────────────────────────────────────────────────────────────────────

def _table_header_row(slide, cols, lefts, top, height, widths, font_size: int = 8):
    for col, left, width in zip(cols, lefts, widths):
        _add_rect(slide, left, top, width, height, ORANGE_DARK, MID_GREY)
        _add_text_box(slide, str(col) if col else "",
                      left + Pt(2), top + Pt(1), width - Pt(4), height - Pt(2),
                      font_size=font_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _table_data_row(slide, cols, lefts, top, height, widths,
                    font_size: int = 8, bg=None, highlight: bool = False):
    fill = ORANGE_LIGHT if highlight else (bg or WHITE)
    txt  = ORANGE_DARK  if highlight else BLACK
    for i, (col, left, width) in enumerate(zip(cols, lefts, widths)):
        _add_rect(slide, left, top, width, height, fill, MID_GREY)
        _add_text_box(slide, str(col) if col else "",
                      left + Pt(2), top + Pt(1), width - Pt(4), height - Pt(2),
                      font_size=font_size, bold=highlight,
                      color=txt, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)


def _variation_color(val_str: str) -> RGBColor:
    s = str(val_str).strip()
    if s.startswith("+") or (s.replace(".", "").replace(",", "").lstrip("0")
                               and not s.startswith("-")
                               and s not in ("—", "0", "0,00%", "0.00%")):
        return GREEN_UP
    if s.startswith("-"):
        return RED_DOWN
    return DARK_GREY


def _summary_oval(slide, label: str, value: str, left, top, width, height,
                  fill=None, text_color=WHITE):
    if fill is None:
        fill = BROWN_KPI
    oval = slide.shapes.add_shape(9, left, top, width, height)
    oval.fill.solid()
    oval.fill.fore_color.rgb = fill
    oval.line.fill.background()
    lbl_h = height * 0.42
    val_h = height - lbl_h
    _add_text_box(slide, label,
                  left + Inches(0.04), top + Inches(0.04),
                  width - Inches(0.08), lbl_h - Inches(0.04),
                  font_size=7, bold=True, color=text_color, align=PP_ALIGN.CENTER)
    _add_text_box(slide, value,
                  left + Inches(0.04), top + lbl_h,
                  width - Inches(0.08), val_h,
                  font_size=9, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Main Generator
# ─────────────────────────────────────────────────────────────────────────────

class BDAReportGenerator:
    """
    Builds the full BDA Daily Report PPTX from template_v1.pptx.

    Top-level data keys
    -------------------
    report_date         str      "30.03.2026"
    kpis                list     [{label, value, variation_str}]
    reembolso_credito   str      "17,62 M Kz"
    liquidez_mn_days    list     ["25/11","26/11","27/11","28/11","01/12"]
    liquidez_mn_rows    list     [{label, values:[5 str]}]
    transacoes_mn_raw   list
    luibor / luibor_d1 / luibor_d2 / luibor_variation  dict
    operacoes_vivas     list
    juros_diario_mn     str
    fluxos_mn_rows      list
    pl_summary          list
    desembolsos_total   float
    reembolsos_pie      list
    liquidez_me_rows    list
    transacoes_me_raw   list
    operacoes_vivas_me  list
    fluxos_me_rows      list
    juros_diario_me     str
    cambial             dict
    cambial_rows        list
    transacoes_bda_rows list
    mercado_rows        list
    bodiva_segment_rows list
    bodiva_total_transacoes str
    bodiva_stocks       dict
    bodiva_operacoes    list
    bodiva_transacoes_valor str
    bodiva_juros_diario str
    carteira_titulos    list
    market_info         dict
    """

    def __init__(self, data=None):
        self.data = data or {}

    # ── Public ────────────────────────────────────────────────────────────────

    def build(self, output_path: str = "output/bda_report.pptx") -> str:
        """Load template, modify slides in-place, save to output_path."""
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

        if not _TEMPLATE.exists():
            raise FileNotFoundError(
                f"Template not found at {_TEMPLATE}. "
                "Place template_v1.pptx in the assets/ directory."
            )

        shutil.copy(str(_TEMPLATE), output_path)
        self.prs = Presentation(output_path)
        slides   = self.prs.slides

        # Slide 0: Cover — update date only
        date_str = self.data.get("report_date", date.today().strftime("%d.%m.%Y"))
        _update_cover_date(slides[0], date_str)

        # Slide 1: Agenda — 100% static, skip

        # Slides 2-10: clear and rebuild with live data
        _pairs = [
            (2,  self._slide_sumario_executivo),
            (3,  self._slide_liquidez_mn_1),
            (4,  self._slide_liquidez_mn_2),
            (5,  self._slide_liquidez_me),
            (6,  self._slide_mercado_cambial),
            (7,  self._slide_bodiva),
            (8,  self._slide_operacoes_bda),
            (9,  self._slide_market_info_1),
            (10, self._slide_market_info_2),
        ]
        for idx, method in _pairs:
            slide = slides[idx]
            if idx != 2:
                _clear_slide(slide)
            method(slide)

        self.prs.save(output_path)
        return output_path

    # ── Chart helpers ─────────────────────────────────────────────────────────

    def _add_pie_charts_mn(self, slide):
        """Slide 5: DESEMBOLSOS + REEMBOLSOS pie charts."""
        import io
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            import matplotlib.patches as mpatches
        except ImportError:
            return

        PIE_COLORS = ["#E8751A", "#C05000", "#A03800", "#FF9A3C", "#7A3000",
                      "#FFC57A", "#5C2D00", "#FFAA55"]

        # ── DESEMBOLSOS ──────────────────────────────────────────────────────
        desembolso_val = float(str(self.data.get("desembolsos_total", 0)
                                   ).replace(",", ".").replace(" ", "") or 0)
        fig_d, ax_d = plt.subplots(figsize=(2.69, 2.53), subplot_kw=dict(aspect="equal"))
        fig_d.patch.set_alpha(0)
        ax_d.set_facecolor("none")
        if desembolso_val <= 0:
            wedge = mpatches.Wedge((0, 0), 1, 0, 360,
                                   facecolor="#E8751A", edgecolor="#C05000", linewidth=1.5)
            ax_d.add_patch(wedge)
            shadow = mpatches.Ellipse((0, -0.12), 2.0, 0.35,
                                      facecolor="#C05000", alpha=0.6, zorder=0)
            ax_d.add_patch(shadow)
            ax_d.set_xlim(-1.3, 1.3); ax_d.set_ylim(-0.6, 1.3)
        else:
            ax_d.pie([desembolso_val], colors=["#E8751A"],
                     wedgeprops={"edgecolor": "#C05000", "linewidth": 1.5})
        ax_d.set_title("DESEMBOLSOS", fontsize=9, fontweight="bold",
                       color="black", pad=4,
                       bbox=dict(boxstyle="round,pad=0.2", fc="#C05000", ec="#C05000"))
        ax_d.axis("off")
        buf_d = io.BytesIO()
        fig_d.savefig(buf_d, format="png", bbox_inches="tight", transparent=True, dpi=150)
        plt.close(fig_d)
        buf_d.seek(0)
        slide.shapes.add_picture(buf_d, Inches(6.21), Inches(3.34), Inches(2.69), Inches(2.53))

        _place_img(slide, IMG_ICON_KZ_CIRCLE, Inches(6.39), Inches(5.98),
                   Inches(0.56), Inches(0.41))

        # ── REEMBOLSOS ───────────────────────────────────────────────────────
        reemb_items = self.data.get("reembolsos_pie", [])
        if not reemb_items:
            fig_r, ax_r = plt.subplots(figsize=(4.09, 2.4), subplot_kw=dict(aspect="equal"))
            fig_r.patch.set_alpha(0); ax_r.set_facecolor("none")
            ax_r.pie([1], colors=["#E8751A"],
                     wedgeprops={"edgecolor": "#C05000", "linewidth": 1.5})
            ax_r.set_title("REEMBOLSOS", fontsize=9, fontweight="bold",
                           color="black", pad=4,
                           bbox=dict(boxstyle="round,pad=0.2", fc="#C05000", ec="#C05000"))
            ax_r.axis("off")
            buf_r = io.BytesIO()
            fig_r.savefig(buf_r, format="png", bbox_inches="tight", transparent=True, dpi=150)
            plt.close(fig_r)
        else:
            labels = [r["label"] for r in reemb_items]
            values = [float(str(r["valor"]).replace(",", ".")) for r in reemb_items]
            colors = PIE_COLORS[:len(values)]
            fig_r, ax_r = plt.subplots(figsize=(4.09, 2.4), subplot_kw=dict(aspect="equal"))
            fig_r.patch.set_alpha(0); ax_r.set_facecolor("none")
            wedges, _, autotexts = ax_r.pie(
                values, labels=None, colors=colors, autopct="%1.0f%%",
                pctdistance=0.75, startangle=90,
                wedgeprops={"edgecolor": "white", "linewidth": 1},
            )
            for at in autotexts:
                at.set_fontsize(8); at.set_color("white"); at.set_fontweight("bold")
            ax_r.legend(wedges, labels, loc="center right",
                        bbox_to_anchor=(1.55, 0.5), fontsize=7, frameon=False)
            ax_r.set_title("REEMBOLSOS", fontsize=9, fontweight="bold",
                           color="black", pad=4,
                           bbox=dict(boxstyle="round,pad=0.2", fc="#C05000", ec="#C05000"))
            ax_r.axis("off")
            buf_r = io.BytesIO()
            fig_r.savefig(buf_r, format="png", bbox_inches="tight", transparent=True, dpi=150)
            plt.close(fig_r)
        buf_r.seek(0)
        slide.shapes.add_picture(buf_r, Inches(8.52), Inches(3.45), Inches(4.09), Inches(2.4))

    def _add_liquidez_me_pie_chart(self, slide):
        """Slide 6: Liquidez ME composition pie chart (upper-right)."""
        import io
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return

        rows = self.data.get("liquidez_me_rows", [])
        segments = []
        for row in rows:
            label = row.get("label", "")
            if "LIQUIDEZ" in label.upper():
                continue
            values = row.get("values", [])
            latest = values[-1] if values else 0
            try:
                numeric = float(str(latest).replace(".", "").replace(",", "."))
            except ValueError:
                numeric = 0.0
            if numeric > 0:
                segments.append((label, numeric))

        if not segments:
            segments = [("Saldo D.O.", 95.0), ("DPs ME", 4.0), ("Colateral", 1.0)]

        labels = [s[0] for s in segments]
        values = [s[1] for s in segments]
        colors = ["#E8751A", "#C05000", "#7A3000", "#FFC57A"]

        fig, ax = plt.subplots(figsize=(4.63, 2.98), subplot_kw=dict(aspect="equal"))
        fig.patch.set_alpha(0); ax.set_facecolor("none")
        wedges, _, autotexts = ax.pie(
            values, colors=colors[:len(values)], startangle=90, counterclock=False,
            autopct=lambda pct: f"{pct:.0f}%" if pct >= 1 else "",
            pctdistance=0.72, wedgeprops={"edgecolor": "white", "linewidth": 1},
        )
        for at in autotexts:
            at.set_fontsize(8); at.set_color("white"); at.set_fontweight("bold")
        ax.legend(wedges, labels, loc="center left", bbox_to_anchor=(0.90, 0.5),
                  fontsize=7, frameon=False)
        ax.set_title("Composição Liquidez ME", fontsize=8, fontweight="bold",
                     color="#333333", pad=4)
        ax.axis("off")
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", transparent=True, dpi=150)
        plt.close(fig)
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(8.571), Inches(0.377), Inches(4.630), Inches(2.981))

    def _add_cambial_charts(self, slide):
        """Slide 7: Posição Cambial bar chart + Taxa USD/AKZ line chart."""
        import io
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return

        cambial = self.data.get("cambial", {})

        # Bar chart
        activos  = float(str(cambial.get("activos_usd",  0)).replace(",", ".") or 0)
        passivos = float(str(cambial.get("passivos_usd", 0)).replace(",", ".") or 0)
        fig1, ax1 = plt.subplots(figsize=(6.52, 2.81))
        fig1.patch.set_alpha(0); ax1.set_facecolor("none")
        bars = ax1.bar(["Activos", "Passivos"], [activos, passivos],
                       color=["#E8751A", "#5C2D00"], edgecolor="white", width=0.5)
        for bar, val in zip(bars, [activos, passivos]):
            if val:
                ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                         f"{val:.1f}", ha="center", va="bottom", fontsize=7, color="#333333")
        ax1.set_title("Posição Cambial (M USD)", fontsize=8, fontweight="bold",
                      color="#333333", pad=6)
        ax1.tick_params(axis="both", labelsize=7)
        ax1.spines[["top", "right"]].set_visible(False)
        ax1.set_ylabel("M USD", fontsize=6, color="#555555")
        ax1.grid(axis="y", alpha=0.2, linewidth=0.6)
        buf1 = io.BytesIO()
        fig1.savefig(buf1, format="png", bbox_inches="tight", transparent=True, dpi=150)
        plt.close(fig1); buf1.seek(0)
        slide.shapes.add_picture(buf1, Inches(6.212), Inches(3.260), Inches(6.521), Inches(2.809))

        # Line chart
        cambial_rows = self.data.get("cambial_rows", [])
        usd_row = next((r for r in cambial_rows
                        if "USD" in r.get("par", "").upper()
                        and "EUR" not in r.get("par", "").upper()), None)
        if usd_row:
            try:
                rates = [float(str(usd_row.get(k, "0")).replace(",", "."))
                         for k in ("anterior2", "anterior", "atual")]
            except ValueError:
                rates = [0, 0, 0]
        else:
            rates = [0, 0, 0]

        fig2, ax2 = plt.subplots(figsize=(6.39, 1.86))
        fig2.patch.set_alpha(0); ax2.set_facecolor("none")
        ax2.plot(["D-2", "D-1", "D"], rates, color="#E8751A", linewidth=2,
                 marker="o", markersize=5, markerfacecolor="#5C2D00")
        ax2.fill_between(["D-2", "D-1", "D"], rates, alpha=0.15, color="#E8751A")
        ax2.set_title("Taxa USD/AKZ", fontsize=8, fontweight="bold", color="#333333", pad=6)
        ax2.tick_params(axis="both", labelsize=7)
        ax2.spines[["top", "right"]].set_visible(False)
        ax2.set_ylabel("AKZ", fontsize=6, color="#555555")
        ax2.grid(axis="y", alpha=0.2, linewidth=0.6)
        buf2 = io.BytesIO()
        fig2.savefig(buf2, format="png", bbox_inches="tight", transparent=True, dpi=150)
        plt.close(fig2); buf2.seek(0)
        slide.shapes.add_picture(buf2, Inches(0.163), Inches(4.812), Inches(6.385), Inches(1.858))

    def _add_bodiva_segmento_pie(self, slide):
        """Slide 8: Segmentado Por Produtos pie chart (upper-right area)."""
        import io
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return

        seg_rows = self.data.get("bodiva_segment_rows", [])
        segments = [(r["segmento"], r["atual"]) for r in seg_rows
                    if "total" not in r.get("segmento", "").lower()]
        if not segments:
            return

        labels = [s[0] for s in segments]
        raw_vals = []
        for _, v in segments:
            try:
                raw_vals.append(float(str(v).replace(".", "").replace(",", ".").replace(" mM Kz", "").replace(" M Kz", "").strip()))
            except ValueError:
                raw_vals.append(0.0)

        if not any(raw_vals):
            return

        colors = ["#E8751A", "#C05000", "#A03800", "#FF9A3C", "#7A3000", "#FFC57A"]
        fig, ax = plt.subplots(figsize=(3.2, 2.6), subplot_kw=dict(aspect="equal"))
        fig.patch.set_alpha(0); ax.set_facecolor("none")
        wedges, _, autotexts = ax.pie(
            raw_vals, labels=None, colors=colors[:len(raw_vals)],
            autopct=lambda pct: f"{pct:.0f}%" if pct >= 3 else "",
            pctdistance=0.75, startangle=90,
            wedgeprops={"edgecolor": "white", "linewidth": 1},
        )
        for at in autotexts:
            at.set_fontsize(7); at.set_color("white"); at.set_fontweight("bold")
        ax.legend(wedges, labels, loc="center left", bbox_to_anchor=(1.0, 0.5),
                  fontsize=6, frameon=False)
        ax.set_title("Segmentos", fontsize=8, fontweight="bold", color="#333333", pad=4)
        ax.axis("off")
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", transparent=True, dpi=150)
        plt.close(fig); buf.seek(0)
        # Position: upper-right, same area as the placeholder TextBox 134 in template
        slide.shapes.add_picture(buf, Inches(9.08), Inches(0.72), Inches(4.05), Inches(3.10))

    # ── Slide methods — each accepts an existing slide ────────────────────────

    def _slide_sumario_executivo(self, slide):
        if any(shape.shape_type == MSO_SHAPE_TYPE.GROUP for shape in slide.shapes):
            self._update_sumario_template(slide)
            return

        _slide_title(slide, "Sumário Executivo")

        # Chevron arrows
        for left_in in (0.07, 0.39):
            sh = slide.shapes.add_shape(13, Inches(left_in), Inches(0.76),
                                        Inches(0.47), Inches(0.38))
            sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE_PRIMARY
            sh.line.fill.background()

        # Central oval KPI
        rc_val = self.data.get("reembolso_credito", "—")
        cx, cy = Inches(5.741), Inches(3.009)
        bw, bh = Inches(1.779), Inches(1.491)
        oval = slide.shapes.add_shape(9, cx, cy, bw, bh)
        oval.fill.solid(); oval.fill.fore_color.rgb = BROWN_KPI
        oval.line.fill.background()
        _add_text_box(slide, rc_val,
                      cx + Inches(0.1), cy + Inches(0.1), bw - Inches(0.2), Inches(0.65),
                      font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, "Reembolsos de Crédito",
                      cx + Inches(0.1), cy + Inches(0.75), bw - Inches(0.2), Inches(0.4),
                      font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Satellite KPI cards — 4 left + 4 right
        kpis = self.data.get("kpis", [
            {"label": "Liquidez MN",           "value": "—", "variation_str": ""},
            {"label": "Liquidez ME",            "value": "—", "variation_str": ""},
            {"label": "Posição Cambial",        "value": "—", "variation_str": ""},
            {"label": "Carteira Títulos",       "value": "—", "variation_str": ""},
            {"label": "Rentabilidade MN",       "value": "—", "variation_str": ""},
            {"label": "Rentabilidade ME",       "value": "—", "variation_str": ""},
            {"label": "Rentabilidade Títulos",  "value": "—", "variation_str": ""},
            {"label": "Reembolsos",             "value": "—", "variation_str": ""},
        ])

        _ROW = [Inches(0.85), Inches(2.18), Inches(3.51), Inches(4.84)]
        _LC  = Inches(0.15)
        _RC  = Inches(10.10)

        kpi_layout = [
            (_LC, _ROW[0], IMG_ICON_S3_LMN,     Inches(0.44), Inches(0.55)),
            (_LC, _ROW[1], IMG_ICON_S3_LME,     Inches(0.43), Inches(0.38)),
            (_LC, _ROW[2], IMG_ICON_S3_CAMBIAL, Inches(0.56), Inches(0.44)),
            (_LC, _ROW[3], IMG_ICON_S3_VBAR,    Inches(0.10), Inches(0.55)),
            (_RC, _ROW[0], IMG_ICON_S3_RENTA,   Inches(0.46), Inches(0.30)),
            (_RC, _ROW[1], IMG_ICON_S3_RENTA,   Inches(0.46), Inches(0.30)),
            (_RC, _ROW[2], IMG_ICON_S3_RENTA,   Inches(0.46), Inches(0.30)),
            (_RC, _ROW[3], IMG_ICON_S3_REEMB,   Inches(0.55), Inches(0.43)),
        ]

        card_w, card_h = Inches(2.65), Inches(1.21)
        for idx, kpi in enumerate(kpis[:8]):
            if idx >= len(kpi_layout):
                break
            vl, vt, icon_path, iw, ih = kpi_layout[idx]
            _place_img(slide, icon_path, vl, vt, iw, ih)
            card_left = vl + Inches(0.50)
            _add_rect(slide, card_left, vt, card_w, card_h, ORANGE_LIGHT, ORANGE_PRIMARY, 1.0)
            _add_text_box(slide, kpi["label"],
                          card_left + Pt(4), vt + Pt(3), card_w - Pt(8), Inches(0.32),
                          font_size=9, color=DARK_GREY, align=PP_ALIGN.LEFT)
            _add_text_box(slide, kpi["value"],
                          card_left + Pt(4), vt + Inches(0.35), card_w - Pt(8), Inches(0.5),
                          font_size=15, bold=True, color=ORANGE_PRIMARY, align=PP_ALIGN.LEFT)
            if kpi.get("variation_str"):
                vc = _variation_color(kpi["variation_str"])
                _add_text_box(slide, kpi["variation_str"],
                              card_left + Pt(4), vt + Inches(0.85), card_w - Pt(8), Inches(0.28),
                              font_size=9, bold=True, color=vc, align=PP_ALIGN.LEFT)

        _place_img(slide, IMG_ICON_S3_REPORT, Inches(10.74), Inches(6.42),
                   Inches(0.35), Inches(0.35))
        _footer(slide)

    def _iter_shapes_recursive(self, shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_shapes_recursive(shape.shapes)

    @staticmethod
    def _norm_label(text: str) -> str:
        return (
            str(text).lower()
            .replace("í", "i")
            .replace("ó", "o")
            .replace("ã", "a")
            .replace("ç", "c")
            .replace("é", "e")
            .replace("ê", "e")
            .replace("ú", "u")
            .replace("á", "a")
            .replace("â", "a")
        )

    @staticmethod
    def _set_shape_text(shape, lines, font_size: int | None = None):
        tf = shape.text_frame
        tf.clear()
        for idx, line in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(line)
            run.font.name = FONT
            if font_size:
                run.font.size = Pt(font_size)
            if idx == 0:
                run.font.bold = True

    def _update_metric_group(self, group, label: str, value: str, variation: str = ""):
        wanted = self._norm_label(label)
        text_shapes = [
            shape for shape in self._iter_shapes_recursive(group.shapes)
            if getattr(shape, "has_text_frame", False)
        ]
        main = next(
            (shape for shape in text_shapes if wanted in self._norm_label(shape.text)),
            None,
        )
        if main:
            self._set_shape_text(main, [label, value], font_size=9)

        var_shape = next(
            (
                shape for shape in text_shapes
                if "face ao dia anterior" in self._norm_label(shape.text)
                and shape is not main
            ),
            None,
        )
        if var_shape:
            self._set_shape_text(
                var_shape,
                [variation or "0,00%", "Face ao dia anterior"],
                font_size=8,
            )

    def _update_sumario_template(self, slide):
        kpi_rows = {
            self._norm_label(item.get("label", "")): item
            for item in self.data.get("kpis", [])
        }

        metrics = [
            ("Liquidez MN", "Liquidez MN"),
            ("Liquidez ME", "Liquidez ME"),
            ("Posição Cambial", "Posição Cambial"),
            ("Carteira Titulos", "Carteira Títulos"),
            ("Rentabilidade MN", "Rentabilidade MN"),
            ("Rentabilidade ME", "Rentabilidade ME"),
            ("Rentabilidade Títulos", "Rentabilidade Títulos"),
            ("Reembolsos", "Reembolsos"),
        ]

        for group in (shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP):
            group_text = " ".join(
                shape.text for shape in self._iter_shapes_recursive(group.shapes)
                if getattr(shape, "has_text_frame", False)
            )
            normalized_group_text = self._norm_label(group_text)
            for match_label, data_label in metrics:
                if self._norm_label(match_label) not in normalized_group_text:
                    continue
                row = kpi_rows.get(self._norm_label(data_label), {})
                value = row.get("value", "—")
                variation = row.get("variation_str", "")
                self._update_metric_group(group, match_label, value, variation)

            if "desembolsos" in normalized_group_text:
                desembolsos = self.data.get("desembolsos_total", "—")
                try:
                    desembolsos_value = f"{float(desembolsos):,.0f} mM Kz".replace(",", ".")
                except (TypeError, ValueError):
                    desembolsos_value = str(desembolsos)
                self._update_metric_group(group, "Desembolsos", desembolsos_value, "")

        rc_val = self.data.get("reembolso_credito", "—")
        for shape in slide.shapes:
            if shape.shape_id == 40 and shape.has_text_frame:
                self._set_shape_text(shape, [rc_val, "Reembolsos de Crédito"], font_size=12)
                break

    # ── Slide 4: Liquidez MN — tables + LUIBOR ───────────────────────────────

    def _slide_liquidez_mn_1(self, slide):
        _slide_title(slide, "LIQUIDEZ – MOEDA NACIONAL")

        days    = self.data.get("liquidez_mn_days", ["D-4", "D-3", "D-2", "D-1", "D"])
        row_h   = Inches(0.28)
        left0   = Inches(0.407)
        tbl_w   = Inches(8.557)
        label_w = Inches(3.0)
        col_w   = (tbl_w - label_w) / 5
        lefts   = [left0] + [left0 + label_w + i * col_w for i in range(5)]
        widths  = [label_w] + [col_w] * 5

        # ── Liquidez MN table ─────────────────────────────────────────────────
        top = Inches(0.70)
        _section_bar(slide, "Liquidez MN  (Em Milhares)", left0, top, tbl_w)
        top += Inches(0.28)
        _table_header_row(slide, [""] + days, lefts, top, row_h, widths)

        lmn_rows = self.data.get("liquidez_mn_rows", [
            {"label": "Posição Reservas Livres BNA", "values": ["—"] * 5},
            {"label": "Posição DO B. Comerciais",    "values": ["—"] * 5},
            {"label": "Posição DP B. Comerciais",    "values": ["—"] * 5},
            {"label": "Posição OMAs",                "values": ["—"] * 5},
            {"label": "LIQUIDEZ BDA",                "values": ["—"] * 5},
        ])
        for i, row in enumerate(lmn_rows):
            top += row_h
            is_total = "LIQUIDEZ" in row["label"].upper()
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide, [row["label"]] + row["values"],
                            lefts, top, row_h, widths, highlight=is_total, bg=bg)

        # ── Transações ────────────────────────────────────────────────────────
        top = Inches(2.204)
        _section_bar(slide, "Transações", left0, top, tbl_w)
        tx_heads  = ["Tipo", "Contraparte", "Taxa", "Montante", "Maturidade", "Juros"]
        tx_w_list = [Inches(1.1), Inches(1.8), Inches(0.9), Inches(1.8), Inches(1.6), Inches(1.357)]
        tx_lefts  = [left0]
        for w in tx_w_list[:-1]:
            tx_lefts.append(tx_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, tx_heads, tx_lefts, top, row_h, tx_w_list)

        raw_tx = self.data.get("transacoes_mn_raw", [])
        for i, row in enumerate(raw_tx):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row.get("tipo", "OMA"), row.get("contraparte", "—"),
                             row.get("taxa", "—"), row.get("montante", "—"),
                             row.get("maturidade", "—"), row.get("juros", "—")],
                            tx_lefts, top, row_h, tx_w_list, bg=bg)
        if not raw_tx:
            top += row_h
            _table_data_row(slide, ["—"] * 6, tx_lefts, top, row_h, tx_w_list)

        # ── Operações Vivas ────────────────────────────────────────────────────
        top = Inches(3.012)
        _section_bar(slide, "Operações Vivas", left0, top, tbl_w)
        op_heads  = ["Tipo", "Contraparte", "Montante", "Taxa", "Residual", "Vencimento", "Juro Diário"]
        n_op      = len(op_heads)
        op_w      = tbl_w / n_op
        op_lefts  = [left0 + i * op_w for i in range(n_op)]
        op_widths = [op_w] * n_op
        top += Inches(0.28)
        _table_header_row(slide, op_heads, op_lefts, top, row_h, op_widths, font_size=7)

        ops = self.data.get("operacoes_vivas", [])
        for i, op in enumerate(ops):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [op.get("tipo", "DP"), op.get("contraparte", "—"),
                             op.get("montante", "—"), op.get("taxa", "—"),
                             str(op.get("residual", "—")), op.get("vencimento", "—"),
                             op.get("juro_diario", "—")],
                            op_lefts, top, row_h, op_widths, font_size=7, bg=bg)
        if not ops:
            top += row_h
            _table_data_row(slide, ["—"] * n_op, op_lefts, top, row_h, op_widths, font_size=7)

        # ── LUIBOR ────────────────────────────────────────────────────────────
        top += row_h + Inches(0.12)
        _section_bar(slide, "Taxas LUIBOR", left0, top, tbl_w)
        tenors    = ["LUIBOR O/N", "LUIBOR 1M", "LUIBOR 3M", "LUIBOR 6M", "LUIBOR 9M", "LUIBOR 12M"]
        lu_heads  = ["Maturidade", "Anterior (D-2)", "Anterior (D-1)", "Actual (D)", "Var (%)"]
        lu_n      = len(lu_heads)
        lu_w      = tbl_w / lu_n
        lu_lefts  = [left0 + i * lu_w for i in range(lu_n)]
        lu_widths = [lu_w] * lu_n
        top += Inches(0.28)
        _table_header_row(slide, lu_heads, lu_lefts, top, row_h, lu_widths)

        luibor     = self.data.get("luibor", {})
        luibor_d1  = self.data.get("luibor_d1", {})
        luibor_d2  = self.data.get("luibor_d2", {})
        luibor_var = self.data.get("luibor_variation", {})
        for i, t in enumerate(tenors):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            rate_d  = luibor.get(t, "—")
            rate_d1 = luibor_d1.get(t, rate_d)
            rate_d2 = luibor_d2.get(t, rate_d1)
            var     = luibor_var.get(t, "—")
            _table_data_row(slide, [t, rate_d2, rate_d1, rate_d, var],
                            lu_lefts, top, row_h, lu_widths, bg=bg)

        # ── Summary ovals ─────────────────────────────────────────────────────
        lmn_rows = self.data.get("liquidez_mn_rows", [])
        total_mn = lmn_rows[-1]["values"][-1] if lmn_rows else "—"
        juros_mn = self.data.get("juros_diario_mn", "—")
        _summary_oval(slide, "Liquidez Total", total_mn,
                      Inches(9.250), Inches(5.450), Inches(1.096), Inches(0.865))
        _summary_oval(slide, "Juros Diário", juros_mn,
                      Inches(10.780), Inches(5.450), Inches(1.096), Inches(0.865))

        _place_img(slide, IMG_ICON_GEAR,       Inches(12.10), Inches(5.55), Inches(0.81), Inches(0.78))
        _place_img(slide, IMG_ICON_CALCULATOR, Inches(12.60), Inches(5.65), Inches(0.57), Inches(0.67))
        _footer(slide)

    # ── Slide 5: Liquidez MN — Cash-flow + P&L ───────────────────────────────

    def _slide_liquidez_mn_2(self, slide):
        _slide_title(slide, "LIQUIDEZ – MOEDA NACIONAL")

        days   = self.data.get("liquidez_mn_days", ["D-4", "D-3", "D-2", "D-1", "D"])
        row_h  = Inches(0.28)
        lbl_w  = Inches(3.8)
        col_w  = Inches(1.72)
        left0  = Inches(0.3)
        lefts  = [left0] + [left0 + lbl_w + i * col_w for i in range(5)]
        widths = [lbl_w] + [col_w] * 5

        def _render_section(label, key, default_labels, start_top):
            _section_bar(slide, label, left0, start_top, SLIDE_W - Inches(0.5))
            t = start_top + Inches(0.28)
            _table_header_row(slide, [""] + days, lefts, t, row_h, widths)
            rows = self.data.get(key, [
                {"label": lbl, "values": ["—"] * 5} for lbl in default_labels
            ])
            for i, row in enumerate(rows):
                t += row_h
                is_total = any(x in row["label"].upper()
                               for x in ("TOTAL", "GAP", "LÍQUIDO", "RESULTADO"))
                bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
                _table_data_row(slide, [row["label"]] + row["values"],
                                lefts, t, row_h, widths, highlight=is_total, bg=bg)
            return t + row_h

        top = Inches(0.78)
        cf_defaults = [
            "Fluxos de Entradas (Cash in flow)",
            "Recebimentos de cupão de títulos",
            "Reembolsos de crédito (+)",
            "Reembolsos de OMA-O/N + Juros",
            "Transferencia a favor conta BNA",
            "Fluxos de Saídas (Cash out flow)",
            "(-) Juros, comissões e outros",
            "Custos Com Pessoal",
            "Fornecimentos e Serviços",
            "Desembolso de crédito (-)",
            "Impostos",
            "Aplicação em OMA",
            "GAP de Liquidez",
        ]
        top = _render_section("Fluxos de Caixa MN", "fluxos_mn_rows", cf_defaults, top)

        # P&L Control - keep this in the lower-left template zone so it never
        # collides with the right-side desembolsos/reembolsos charts.
        pl_left = Inches(0.407)
        pl_top = Inches(5.600)
        pl_width = Inches(5.770)
        pl_row_h = Inches(0.20)
        _section_bar(slide, "P&L Control", pl_left, pl_top, pl_width, height=Inches(0.24))
        pl_heads = ["Categoria", "Nº Operações", "Montante"]
        pl_w = [pl_width * r for r in (0.50, 0.22, 0.28)]
        pl_lefts = [pl_left, pl_left + pl_w[0], pl_left + pl_w[0] + pl_w[1]]
        top = pl_top + Inches(0.24)
        _table_header_row(slide, pl_heads, pl_lefts, top, pl_row_h, pl_w, font_size=7)

        pl_summary = self.data.get("pl_summary", [
            {"label": "Reembolso de Crédito", "n_ops": "—", "montante": "—"},
            {"label": "Fornecedores",          "n_ops": "—", "montante": "—"},
            {"label": "Desembolso de Crédito", "n_ops": "—", "montante": "—"},
        ])
        for i, row in enumerate(pl_summary):
            top += pl_row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["label"], str(row["n_ops"]), str(row["montante"])],
                            pl_lefts, top, pl_row_h, pl_w, font_size=7, bg=bg)

        self._add_pie_charts_mn(slide)
        _footer(slide)

    # ── Slide 6: Liquidez ME ──────────────────────────────────────────────────

    def _slide_liquidez_me(self, slide):
        _slide_title(slide, "LIQUIDEZ – MOEDA ESTRANGEIRA")

        days  = self.data.get("liquidez_mn_days", ["D-4", "D-3", "D-2", "D-1", "D"])
        row_h = Inches(0.25)
        fs    = 7

        left0  = Inches(0.396)
        tbl_w  = Inches(7.784)
        lbl_w  = Inches(2.3)
        col_w  = (tbl_w - lbl_w) / 5
        lefts  = [left0] + [left0 + lbl_w + i * col_w for i in range(5)]
        widths = [lbl_w] + [col_w] * 5

        # Liquidez ME
        top = Inches(0.677)
        _section_bar(slide, "Liquidez ME  (M USD)", left0, top, tbl_w)
        top += Inches(0.25)
        _table_header_row(slide, [""] + days, lefts, top, row_h, widths, font_size=fs)

        lme_rows = self.data.get("liquidez_me_rows", [
            {"label": "SALDO D.O Estrangeiros", "values": ["—"] * 5},
            {"label": "DPs ME",                 "values": ["—"] * 5},
            {"label": "COLATERAL CDI",           "values": ["—"] * 5},
            {"label": "LIQUIDEZ BDA",            "values": ["—"] * 5},
        ])
        for i, row in enumerate(lme_rows):
            top += row_h
            is_total = "LIQUIDEZ" in row["label"].upper()
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide, [row["label"]] + row["values"],
                            lefts, top, row_h, widths, font_size=fs,
                            highlight=is_total, bg=bg)

        # Transações ME
        top = Inches(1.838)
        _section_bar(slide, "Transações ME", left0, top, tbl_w)
        tx_me_heads = ["Tipo", "Moeda", "Contraparte", "Taxa", "Montante", "Maturidade", "Juros"]
        n_tx  = len(tx_me_heads)
        tx_w  = tbl_w / n_tx
        tx_lefts  = [left0 + i * tx_w for i in range(n_tx)]
        tx_widths = [tx_w] * n_tx
        top += Inches(0.25)
        _table_header_row(slide, tx_me_heads, tx_lefts, top, row_h, tx_widths, font_size=fs)

        tx_me = self.data.get("transacoes_me_raw", [])
        for i, row in enumerate(tx_me):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row.get(k, "—") for k in
                             ("tipo", "moeda", "contraparte", "taxa", "montante", "maturidade", "juros")],
                            tx_lefts, top, row_h, tx_widths, font_size=fs, bg=bg)
        if not tx_me:
            top += row_h
            _table_data_row(slide, ["—"] * n_tx, tx_lefts, top, row_h, tx_widths, font_size=fs)

        # Operações Vivas ME
        top = Inches(2.819)
        _section_bar(slide, "Operações Vivas ME", left0, top, tbl_w)
        op_me_heads = ["Contraparte", "Montante", "Taxa", "Residual (Dias)", "Vencimento", "Juro Diário"]
        n_op = len(op_me_heads)
        op_w = tbl_w / n_op
        op_lefts  = [left0 + i * op_w for i in range(n_op)]
        op_widths = [op_w] * n_op
        top += Inches(0.25)
        _table_header_row(slide, op_me_heads, op_lefts, top, row_h, op_widths, font_size=fs)

        ops_me = self.data.get("operacoes_vivas_me", [])
        for i, op in enumerate(ops_me):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [op.get(k, "—") for k in
                             ("contraparte", "montante", "taxa", "residual", "vencimento", "juro_diario")],
                            op_lefts, top, row_h, op_widths, font_size=fs, bg=bg)
        if not ops_me:
            top += row_h
            _table_data_row(slide, ["—"] * n_op, op_lefts, top, row_h, op_widths, font_size=fs)

        # Fluxos de Caixa ME (right side)
        R_left  = Inches(8.23)
        R_width = Inches(5.021)
        fx_lbl_w = Inches(1.8)
        fx_col_w = (R_width - fx_lbl_w) / 5
        fx_lefts  = [R_left] + [R_left + fx_lbl_w + i * fx_col_w for i in range(5)]
        fx_widths = [fx_lbl_w] + [fx_col_w] * 5

        r_top = Inches(3.078)
        _section_bar(slide, "Fluxos de Caixa ME", R_left, r_top, R_width)
        r_top += Inches(0.25)
        _table_header_row(slide, [""] + days, fx_lefts, r_top, row_h, fx_widths, font_size=fs)

        fluxos_me = self.data.get("fluxos_me_rows", [
            {"label": "Cash in flow",      "values": ["—"] * 5},
            {"label": "Outros receb.",     "values": ["—"] * 5},
            {"label": "Remb. DP + Juros",  "values": ["—"] * 5},
            {"label": "Cash out flow",     "values": ["—"] * 5},
            {"label": "Aplicação DP ME",   "values": ["—"] * 5},
            {"label": "GAP de Liquidez",   "values": ["—"] * 5},
        ])
        for i, row in enumerate(fluxos_me):
            r_top += row_h
            is_total = "GAP" in row["label"].upper() or "TOTAL" in row["label"].upper()
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide, [row["label"]] + row["values"],
                            fx_lefts, r_top, row_h, fx_widths, font_size=fs,
                            highlight=is_total, bg=bg)

        self._add_liquidez_me_pie_chart(slide)

        lme_data = self.data.get("liquidez_me_rows", [])
        total_me = lme_data[-1]["values"][-1] if lme_data else "—"
        juros_me = self.data.get("juros_diario_me", "—")
        _summary_oval(slide, "Liquidez Total", total_me,
                      Inches(2.817), Inches(5.260), Inches(1.176), Inches(0.941))
        _summary_oval(slide, "Juros Diário", juros_me,
                      Inches(4.350), Inches(5.260), Inches(1.175), Inches(0.941))
        _place_img(slide, IMG_ICON_CALCULATOR, Inches(5.884), Inches(5.607),
                   Inches(0.615), Inches(0.712))
        _footer(slide)

    # ── Slide 7: Mercado Cambial ──────────────────────────────────────────────

    def _slide_mercado_cambial(self, slide):
        _slide_title(slide, "MERCADO CAMBIAL")

        cambial  = self.data.get("cambial", {})
        row_h    = Inches(0.28)
        left0    = Inches(0.396)
        table_w  = Inches(6.112)

        # Cambiais rates
        top = Inches(0.705)
        _section_bar(slide, "Cambiais", left0, top, table_w)
        c_heads = ["Par", "Anterior (D-2)", "Anterior (D-1)", "Actual (D)", "(%)"]
        c_w     = [Inches(1.4), Inches(1.178), Inches(1.178), Inches(1.178), Inches(1.178)]
        c_lefts = [left0]
        for w in c_w[:-1]:
            c_lefts.append(c_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, c_heads, c_lefts, top, row_h, c_w)

        cambial_rows = self.data.get("cambial_rows", [
            {"par": "USD/AKZ", "anterior2": "—", "anterior": "—", "atual": "—", "variacao": "—"},
            {"par": "EUR/AKZ", "anterior2": "—", "anterior": "—", "atual": "—", "variacao": "—"},
            {"par": "EUR/USD", "anterior2": "—", "anterior": "—", "atual": "—", "variacao": "—"},
        ])
        for i, row in enumerate(cambial_rows):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["par"], row.get("anterior2", "—"), row.get("anterior", "—"),
                             row.get("atual", "—"), row.get("variacao", "—")],
                            c_lefts, top, row_h, c_w, bg=bg)

        # Transações BDA
        top = Inches(1.865)
        _section_bar(slide, "Transações BDA", left0, top, table_w)
        tb_heads = ["C/V", "Par de moeda", "Montante Debt", "Câmbio", "P/L AKZ"]
        tb_w     = [Inches(0.7), Inches(1.412), Inches(1.5), Inches(1.3), Inches(1.2)]
        tb_lefts = [left0]
        for w in tb_w[:-1]:
            tb_lefts.append(tb_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, tb_heads, tb_lefts, top, row_h, tb_w, font_size=7)

        bda_tx = self.data.get("transacoes_bda_rows", [])
        for i, row in enumerate(bda_tx):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row.get(k, "—") for k in ("cv", "par", "montante", "cambio", "pl")],
                            tb_lefts, top, row_h, tb_w, font_size=7, bg=bg)
        if not bda_tx:
            top += row_h
            _table_data_row(slide, ["—"] * 5, tb_lefts, top, row_h, tb_w, font_size=7)

        # Transações do Mercado
        top = Inches(3.047)
        _section_bar(slide, "Transações do Mercado", left0, top, Inches(5.417))
        tm_heads = ["Liquidação", "Montante USD", "Mínimo", "Máximo"]
        tm_w     = [Inches(1.2), Inches(1.539), Inches(1.339), Inches(1.339)]
        tm_lefts = [left0]
        for w in tm_w[:-1]:
            tm_lefts.append(tm_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, tm_heads, tm_lefts, top, row_h, tm_w)

        mercado_rows = self.data.get("mercado_rows", [
            {"label": "T+0", "montante": "—", "min": "—", "max": "—"},
            {"label": "T+1", "montante": "—", "min": "—", "max": "—"},
            {"label": "T+2", "montante": "—", "min": "—", "max": "—"},
        ])
        for i, row in enumerate(mercado_rows):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row.get("label", "—"), row.get("montante", "—"),
                             row.get("min", "—"), row.get("max", "—")],
                            tm_lefts, top, row_h, tm_w, bg=bg)

        # KPI ovals
        _summary_oval(slide, "Transações (USD)", cambial.get("vol_total_usd", "—"),
                      Inches(8.502), Inches(1.422), Inches(1.241), Inches(0.976))
        _summary_oval(slide, "Posição Cambial (Kz)", cambial.get("posicao_cambial", "—"),
                      Inches(10.308), Inches(1.454), Inches(1.241), Inches(0.954))

        # Chart title bars
        _add_rect(slide, Inches(8.131), Inches(2.957), Inches(3.256), Inches(0.303),
                  ORANGE_LIGHT, ORANGE_PRIMARY, 0.8)
        _add_text_box(slide, "Posição Cambial",
                      Inches(8.131), Inches(2.982), Inches(3.256), Inches(0.220),
                      font_size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        _add_rect(slide, Inches(2.033), Inches(4.526), Inches(3.256), Inches(0.303),
                  ORANGE_LIGHT, ORANGE_PRIMARY, 0.8)
        _add_text_box(slide, "Taxa de Cambio",
                      Inches(2.033), Inches(4.551), Inches(3.256), Inches(0.220),
                      font_size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

        self._add_cambial_charts(slide)

        _place_img(slide, IMG_ICON_FX_EXCHANGE,  Inches(12.089), Inches(0.667),
                   Inches(0.473), Inches(0.467))
        _place_img(slide, IMG_ICON_REPORT_MONEY, Inches(6.554),  Inches(2.453),
                   Inches(0.543), Inches(0.536))
        _footer(slide)

    # ── Slide 8: Mercado de Capitais – BODIVA ─────────────────────────────────

    def _slide_bodiva(self, slide):
        _slide_title(slide, "MERCADO DE CAPITAIS")

        row_h = Inches(0.28)
        left0 = Inches(0.396)

        # Segmentado por Produtos
        top = Inches(0.721)
        seg_tbl_w = Inches(7.494)
        _section_bar(slide, "Segmentado Por Produtos", left0, top, seg_tbl_w)
        sp_heads = ["Segmento", "Anterior", "Actual", "(%)"]
        sp_w     = [Inches(3.0), Inches(1.498), Inches(1.498), Inches(1.498)]
        sp_lefts = [left0]
        for w in sp_w[:-1]:
            sp_lefts.append(sp_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, sp_heads, sp_lefts, top, row_h, sp_w)

        seg_rows = self.data.get("bodiva_segment_rows", [
            {"segmento": "Obrigações De Tesouro",    "anterior": "—", "atual": "—", "variacao": "—"},
            {"segmento": "Bilhetes Do Tesouro",       "anterior": "—", "atual": "—", "variacao": "—"},
            {"segmento": "Obrigações Privadas",       "anterior": "—", "atual": "—", "variacao": "—"},
            {"segmento": "Unidades De Participações", "anterior": "—", "atual": "—", "variacao": "—"},
            {"segmento": "Acções",                    "anterior": "—", "atual": "—", "variacao": "—"},
            {"segmento": "Repos",                     "anterior": "—", "atual": "—", "variacao": "—"},
            {"segmento": "Total",                     "anterior": "—", "atual": "—", "variacao": "—"},
        ])
        for i, row in enumerate(seg_rows):
            top += row_h
            is_total = "TOTAL" in row["segmento"].upper()
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["segmento"], row["anterior"], row["atual"], row["variacao"]],
                            sp_lefts, top, row_h, sp_w, highlight=is_total, bg=bg)

        # KPI oval
        total_atual = self.data.get("bodiva_total_transacoes", "—")
        _summary_oval(slide, "Kz  Transações", total_atual,
                      Inches(4.721), Inches(3.159), Inches(1.473), Inches(1.058))

        # Segmento pie chart (right side)
        self._add_bodiva_segmento_pie(slide)

        # Mercado de Bolsas de Acções
        top = Inches(4.622)
        _section_bar(slide, "Mercado de Bolsas de Acções", left0, top, SLIDE_W - Inches(0.5))
        stk_heads = ["Código", "Vol. Transacc.", "Preço Anterior", "Preço Actual", "Variação", "Cap. Bolsista"]
        stk_w     = [Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.0), Inches(1.8), Inches(3.73)]
        stk_lefts = [left0]
        for w in stk_w[:-1]:
            stk_lefts.append(stk_lefts[-1] + stk_w[len(stk_lefts) - 1])
        top += Inches(0.28)
        _table_header_row(slide, stk_heads, stk_lefts, top, row_h, stk_w)

        stocks = self.data.get("bodiva_stocks", {})
        for i, (code, info) in enumerate(stocks.items()):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            chg = info.get("change_pct")
            chg_str = f"{chg:+.2f}%" if isinstance(chg, (int, float)) else str(chg or "—")
            _table_data_row(slide,
                            [code, str(info.get("volume", "—")),
                             str(info.get("previous", "—")), str(info.get("current", "—")),
                             chg_str, str(info.get("cap_bolsista", "—"))],
                            stk_lefts, top, row_h, stk_w, bg=bg)
        if not stocks:
            top += row_h
            _table_data_row(slide, ["Dados não disponíveis (BODIVA)"] + ["—"] * 5,
                            stk_lefts, top, row_h, stk_w)

        _footer(slide)

    # ── Slide 9: Operações BDA ────────────────────────────────────────────────

    def _slide_operacoes_bda(self, slide):
        _slide_title(slide, "MERCADO DE CAPITAIS – OPERAÇÕES BDA")

        row_h = Inches(0.26)
        left0 = Inches(0.3)

        # Transações
        top = Inches(0.78)
        _section_bar(slide, "Transações", left0, top, SLIDE_W - Inches(0.5))
        tx_heads  = ["Tipo de Operação", "Data Contrat.", "C/V", "Preço", "Quantidades", "Montante"]
        n_tx = len(tx_heads)
        tx_w = (SLIDE_W - Inches(0.5)) / n_tx
        tx_lefts  = [left0 + i * tx_w for i in range(n_tx)]
        tx_widths = [tx_w] * n_tx
        top += Inches(0.28)
        _table_header_row(slide, tx_heads, tx_lefts, top, row_h, tx_widths, font_size=7)

        ops = self.data.get("bodiva_operacoes", [])
        for i, row in enumerate(ops):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row.get(k, "—") for k in
                             ("tipo", "data", "cv", "preco", "quantidade", "montante")],
                            tx_lefts, top, row_h, tx_widths, font_size=7, bg=bg)
        if not ops:
            top += row_h
            _table_data_row(slide, ["—"] * n_tx, tx_lefts, top, row_h, tx_widths, font_size=7)

        # Carteira de Títulos
        top = Inches(1.85)
        carteira  = self.data.get("carteira_titulos", [])
        custo_rows = [r for r in carteira if "CUSTO" in r.get("carteira", "").upper()]
        justo_rows = [r for r in carteira if "JUSTO" in r.get("carteira", "").upper()
                      or "VALOR" in r.get("carteira", "").upper()]
        if not custo_rows and not justo_rows:
            mid = max(1, len(carteira) // 2)
            custo_rows, justo_rows = carteira[:mid], carteira[mid:]

        ct_heads  = ["Cód.", "Qtd D-1", "Qtd D", "Nominal", "Taxa", "Montante", "J. Anual", "J. Diário"]
        n_ct      = len(ct_heads)
        ct_w_each = (SLIDE_W - Inches(0.5)) / n_ct
        ct_lefts  = [left0 + i * ct_w_each for i in range(n_ct)]
        ct_widths = [ct_w_each] * n_ct

        def _render_carteira_section(label, rows, start_top):
            _section_bar(slide, label, left0, start_top, SLIDE_W - Inches(0.5))
            t = start_top + Inches(0.24)
            _table_header_row(slide, ct_heads, ct_lefts, t, row_h, ct_widths, font_size=6)
            if rows:
                for j, row in enumerate(rows):
                    t += row_h
                    is_total = row.get("cod", "").upper() == "TOTAL"
                    bg = ORANGE_LIGHT if j % 2 == 0 else WHITE
                    _table_data_row(slide,
                                    [row.get(k, "—") for k in
                                     ("cod", "qty_d1", "qty_d", "nominal", "taxa",
                                      "montante", "juros_anual", "juro_diario")],
                                    ct_lefts, t, row_h, ct_widths, font_size=6,
                                    highlight=is_total, bg=bg)
            else:
                t += row_h
                _table_data_row(slide, ["—"] * n_ct, ct_lefts, t, row_h, ct_widths, font_size=6)
            return t + row_h

        top = _render_carteira_section("Custo Amortizado", custo_rows, top)
        top += Inches(0.08)
        _render_carteira_section("Justo Valor", justo_rows, top)

        # KPI ovals
        tx_kpi_val = self.data.get("bodiva_transacoes_valor", "—")
        jd_kpi_val = self.data.get("bodiva_juros_diario", "—")
        _summary_oval(slide, "Kz  Transações",   tx_kpi_val,
                      Inches(4.50), Inches(5.85), Inches(1.30), Inches(0.88))
        _summary_oval(slide, "Kz  Juros Diário", jd_kpi_val,
                      Inches(6.00), Inches(5.85), Inches(1.30), Inches(0.88))
        _footer(slide)

    # ── Slide 10: Informação de Mercados (1/2) ────────────────────────────────

    def _slide_market_info_1(self, slide):
        _slide_title(slide, "INFORMAÇÃO DE MERCADOS")

        market = self.data.get("market_info", {})
        row_h  = Inches(0.28)
        left0  = Inches(0.396)
        tbl_w  = Inches(6.0)
        com_x  = Inches(7.134)
        com_w  = Inches(5.702)

        # Capital Markets
        top = Inches(0.693)
        _section_bar(slide, "Capital Markets", left0, top, tbl_w)
        cm_heads = ["Índice", "Anterior", "Actual", "(%)"]
        cm_w     = [Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.2)]
        cm_lefts = [left0]
        for w in cm_w[:-1]:
            cm_lefts.append(cm_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, cm_heads, cm_lefts, top, row_h, cm_w)

        cm_rows = market.get("capital_markets", [
            {"indice": "S&P500",            "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "Dow Jones",          "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "NASDAQ",             "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "NIKKEI 225",         "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "IBOVESPA",           "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "EUROSTOX",           "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "Bolsa de Londres",   "anterior": "—", "atual": "—", "variacao": "—"},
            {"indice": "PSI 20",             "anterior": "—", "atual": "—", "variacao": "—"},
        ])
        for i, row in enumerate(cm_rows):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["indice"], row["anterior"], row["atual"], row["variacao"]],
                            cm_lefts, top, row_h, cm_w, bg=bg)

        cm_comment = market.get("cm_commentary", "")
        if cm_comment:
            _add_rect(slide, com_x, Inches(0.662), com_w, Inches(2.077),
                      ORANGE_LIGHT, ORANGE_PRIMARY, 0.8)
            _add_text_box(slide, cm_comment,
                          com_x + Pt(6), Inches(0.69), com_w - Pt(12), Inches(2.0),
                          font_size=8, color=BLACK, word_wrap=True)

        # "Nota" tag
        _add_rect(slide, Inches(12.340), Inches(3.327), Inches(0.598), Inches(0.315),
                  ORANGE_PRIMARY)
        _add_text_box(slide, "Nota", Inches(12.344), Inches(3.331),
                      Inches(0.590), Inches(0.307),
                      font_size=8, bold=True, italic=True,
                      color=WHITE, align=PP_ALIGN.CENTER)

        # Criptomoedas
        top = Inches(4.039)
        _section_bar(slide, "Criptomoedas", left0, top, tbl_w)
        cr_heads = ["Moeda", "Anterior", "Actual", "(%)"]
        cr_w     = [Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.2)]
        cr_lefts = [left0]
        for w in cr_w[:-1]:
            cr_lefts.append(cr_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, cr_heads, cr_lefts, top, row_h, cr_w)

        cr_rows = market.get("crypto", [
            {"moeda": "BITCOIN (BTC)",  "anterior": "—", "atual": "—", "variacao": "—"},
            {"moeda": "ETHEREUM (ETH)", "anterior": "—", "atual": "—", "variacao": "—"},
            {"moeda": "XRP (XRP)",      "anterior": "—", "atual": "—", "variacao": "—"},
        ])
        for i, row in enumerate(cr_rows):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["moeda"], row["anterior"], row["atual"], row["variacao"]],
                            cr_lefts, top, row_h, cr_w, bg=bg)

        cr_comment = market.get("crypto_commentary", "")
        if cr_comment:
            _add_rect(slide, Inches(7.083), Inches(4.068), Inches(5.804), Inches(1.386),
                      ORANGE_LIGHT, ORANGE_PRIMARY, 0.8)
            _add_text_box(slide, cr_comment,
                          Inches(7.10), Inches(4.09), Inches(5.760), Inches(1.340),
                          font_size=8, color=BLACK, word_wrap=True)

        _place_img(slide, IMG_ICON_REPORT_MONEY, Inches(6.43), Inches(0.75),
                   Inches(0.39), Inches(0.39))
        _footer(slide)

    # ── Slide 11: Informação de Mercados (2/2) ────────────────────────────────

    def _slide_market_info_2(self, slide):
        _slide_title(slide, "INFORMAÇÃO DE MERCADOS")

        market = self.data.get("market_info", {})
        row_h  = Inches(0.28)
        left0  = Inches(0.396)
        tbl_w  = Inches(6.0)

        # Commodities
        top = Inches(0.675)
        _section_bar(slide, "Commodities", left0, top, tbl_w)
        cmd_heads = ["Commodity", "Anterior", "Actual", "(%)"]
        cmd_w     = [Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.2)]
        cmd_lefts = [left0]
        for w in cmd_w[:-1]:
            cmd_lefts.append(cmd_lefts[-1] + w)
        top += Inches(0.28)
        _table_header_row(slide, cmd_heads, cmd_lefts, top, row_h, cmd_w)

        cmd_rows = market.get("commodities", [
            {"nome": "PETRÓLEO (BRENT)",       "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "MILHO (USD/BU)",          "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "SOJA (USD/BU)",           "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "TRIGO (USD/LBS)",         "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "CAFÉ (USD/LBS)",          "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "AÇÚCAR (USD/LBS)",        "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "ÓLEO DE PALMA (USD/LBS)", "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "ALGODÃO (USD/LBS)",       "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "BANANA (USD/LBS)",        "anterior": "—", "atual": "—", "variacao": "—"},
        ])
        for i, row in enumerate(cmd_rows):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["nome"], row["anterior"], row["atual"], row["variacao"]],
                            cmd_lefts, top, row_h, cmd_w, bg=bg)

        cmd_comment = market.get("commodities_nota", market.get("commodities_commentary", ""))
        if cmd_comment:
            _add_rect(slide, Inches(7.052), Inches(0.723), Inches(5.572), Inches(1.572),
                      ORANGE_LIGHT, ORANGE_PRIMARY, 0.8)
            _add_text_box(slide, cmd_comment,
                          Inches(7.08), Inches(0.75), Inches(5.53), Inches(1.52),
                          font_size=8, color=BLACK, word_wrap=True)

        # Minerais
        top += row_h + Inches(0.18)
        _section_bar(slide, "Minerais", left0, top, tbl_w)
        min_heads = ["Mineral", "Anterior", "Actual", "(%)"]
        top += Inches(0.28)
        _table_header_row(slide, min_heads, cmd_lefts, top, row_h, cmd_w)

        min_rows = market.get("minerais", [
            {"nome": "OURO",     "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "FERRO",    "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "COBRE",    "anterior": "—", "atual": "—", "variacao": "—"},
            {"nome": "MANGANÊS", "anterior": "—", "atual": "—", "variacao": "—"},
        ])
        for i, row in enumerate(min_rows):
            top += row_h
            bg = ORANGE_LIGHT if i % 2 == 0 else WHITE
            _table_data_row(slide,
                            [row["nome"], row["anterior"], row["atual"], row["variacao"]],
                            cmd_lefts, top, row_h, cmd_w, bg=bg)

        min_comment = market.get("minerais_commentary", "")
        if min_comment:
            _add_rect(slide, Inches(7.187), Inches(3.792), Inches(5.465), Inches(1.413),
                      ORANGE_LIGHT, ORANGE_PRIMARY, 0.8)
            _add_text_box(slide, min_comment,
                          Inches(7.22), Inches(3.82), Inches(5.42), Inches(1.37),
                          font_size=8, color=BLACK, word_wrap=True)

        _place_img(slide, IMG_ICON_REPORT_MONEY, Inches(6.28), Inches(0.90),
                   Inches(0.39), Inches(0.39))
        _footer(slide)
