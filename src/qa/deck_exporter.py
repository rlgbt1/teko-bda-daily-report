"""
qa/deck_exporter.py — Export a generated PPTX to a structured JSON summary.

The exported dict is used by template_checker.py for deterministic checks
and passed to Gemini for template compliance QA.

Output shape:
{
    "slide_count": 11,
    "slides": [
        {
            "index": 1,              # 1-based
            "title": "...",          # first large text found, or ""
            "shape_count": 10,
            "text_shape_count": 8,
            "placeholder_count": 2,  # shapes with only N/A or — values
            "footer_present": True,
            "texts": ["...", ...],   # all non-empty text strings (truncated)
            "table_row_counts": [5, 3],  # one entry per table shape
            "commentary_length": 0,  # chars in commentary block (if detected)
            "na_cell_count": 0,      # table cells that are N/A / —
        },
        ...
    ]
}
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from src.utils.logger import get_logger

logger = get_logger(__name__)

NA_PAT = re.compile(r"^(N/A|—|-{1,3}|n/a|na|null)$", re.I)

# Known commentary slide indices (1-based) and keywords
COMMENTARY_SLIDE_INDICES = {10, 11}
COMMENTARY_KEYWORDS = re.compile(
    r"(mercados|commodities|criptomoedas|minerais|tendência|tendencia|variação)",
    re.I,
)


def export_deck(pptx_path: str | Path) -> dict[str, Any]:
    """
    Read a .pptx file and return a structured summary dict.

    Returns an error dict if the file cannot be read.
    """
    path = Path(pptx_path)
    if not path.exists():
        logger.error("deck_exporter: file not found: %s", path)
        return {"error": f"File not found: {path}", "slide_count": 0, "slides": []}

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        logger.error("python-pptx not installed")
        return {"error": "python-pptx not installed", "slide_count": 0, "slides": []}

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        logger.error("Failed to open PPTX %s: %s", path, exc)
        return {"error": str(exc), "slide_count": 0, "slides": []}

    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_info = _export_slide(idx, slide)
        slides_data.append(slide_info)

    result = {
        "slide_count": len(slides_data),
        "slides": slides_data,
        "file": str(path.name),
    }
    logger.debug("deck_exporter: exported %d slides from %s", len(slides_data), path.name)
    return result


def _export_slide(idx: int, slide) -> dict[str, Any]:
    texts: list[str] = []
    table_row_counts: list[int] = []
    na_cell_count = 0
    placeholder_count = 0
    footer_present = False
    commentary_length = 0

    for shape in slide.shapes:
        # Text frames
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line[:200])   # truncate very long lines

        # Tables
        if shape.has_table:
            table = shape.table
            row_count = len(table.rows)
            table_row_counts.append(row_count)
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if NA_PAT.match(cell_text):
                        na_cell_count += 1

        # Detect footer (small text near bottom of slide containing "UMA VISÃO" or address)
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            if "UMA VISÃO" in full_text.upper() or "EDIFÍCIO BDA" in full_text.upper():
                footer_present = True

    # Count shapes with only NA/placeholder content
    for shape in slide.shapes:
        if shape.has_text_frame:
            all_lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if all_lines and all(NA_PAT.match(t) for t in all_lines):
                placeholder_count += 1

    # Detect commentary (long text block on commentary slides)
    if idx in COMMENTARY_SLIDE_INDICES:
        for t in texts:
            if COMMENTARY_KEYWORDS.search(t) or len(t) > 100:
                commentary_length = max(commentary_length, len(t))

    # Best guess at slide title: first short text that's not a date/address
    title = ""
    for t in texts:
        if 4 < len(t) < 100 and not re.match(r"^\d{2}[/\-]\d{2}", t):
            title = t
            break

    return {
        "index":            idx,
        "title":            title,
        "shape_count":      len(list(slide.shapes)),
        "text_shape_count": sum(1 for s in slide.shapes if s.has_text_frame),
        "placeholder_count": placeholder_count,
        "footer_present":   footer_present,
        "texts":            texts[:30],          # cap at 30 lines per slide
        "table_row_counts": table_row_counts,
        "commentary_length": commentary_length,
        "na_cell_count":    na_cell_count,
    }
